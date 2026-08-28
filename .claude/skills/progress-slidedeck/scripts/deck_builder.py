"""Reusable slide builders for SSP-Plus progress-update decks.

Matches the visual style of the decks in docs/progress-updates/ (e.g.
progress_update_2026-07-16.pptx, progress_update_2026-07-23.pptx) — see
../reference/style.md for the design spec this was derived from. One function per recurring slide type; a new
deck should be a short content script that calls these, not hand-built
shapes/XML.

Usage sketch:

    from deck_builder import *

    prs = new_deck()
    add_title_slide(prs, "Title\nLine Two", "Scope A   ·   Scope B", "August 7, 2026")
    add_overview_slide(prs, "From X to Y", "One or two sentence narrative.",
                        [("Connect", "wire the GUI to the backend"),
                         ("Verify", "confirm it under load"),
                         ("Scale", "generalize to N files")])
    add_flow_slide(prs, "Session 1 — How It Works",
                    [("OTP Screen", "user types the 6-digit code"),
                     ("verify_otp_for_source()", "new: scans pending sessions"),
                     ("file_browser", "real uploaded file loaded")])
    add_built_slide(prs, "Session 1 — What Was Built",
                     [("Session Lookup by Code", ["bullet one", "bullet two"]),
                      ("Email Poller Started", ["bullet one", "bullet two"])])
    add_stats_slide(prs, "Tested and Working", "75", "automated tests, all passing offline",
                     "Covered Offline", ["item", "item"],
                     "Needs the Real Thing", ["item", "item"])
    add_next_steps_slide(prs, "What's Next",
                          [("Verify on the Kiosk", "one sentence description"),
                           ("Add Browser Coverage", "one sentence description")])
    save(prs, "docs/progress-updates/progress_update_2026-08-07.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FONT = "Calibri"

NAVY = RGBColor(0x1A, 0x1A, 0x2E)
DEEP_BLUE = RGBColor(0x1E, 0x27, 0x61)
MID_BLUE = RGBColor(0x2A, 0x35, 0x70)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BLUE = RGBColor(0xE8, 0xEE, 0xFB)
LIGHT_BLUE_2 = RGBColor(0xCA, 0xDC, 0xFC)
AMBER = RGBColor(0xF2, 0xA9, 0x3B)
AMBER_TINT = RGBColor(0xFC, 0xEB, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def save(prs, path):
    prs.save(path)


# ---------------------------------------------------------------- internals

def _blank_slide(prs, bg_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    el = bg._element
    el.getparent().remove(el)
    slide.shapes._spTree.insert(2, el)
    return slide


def _textbox(slide, left, top, width, height, lines, *, size=18, color=WHITE,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.0, space_after=0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return tb


def _rounded_box(slide, left, top, width, height, fill_color, *, radius=0.08):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.fill.background()
    box.shadow.inherit = False
    try:
        box.adjustments[0] = radius
    except (IndexError, ValueError):
        pass
    return box


def _down_arrow(slide, cx, top, size, color):
    left = cx - size // 2
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, size, int(size * 1.1))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    return arrow


def _right_arrow(slide, left, cy, width, color):
    height = Inches(0.35)
    top = cy - height // 2
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    return arrow


def _bullets_into(tf, bullets, *, size=15, color=NAVY, font=FONT, space_after=6):
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = f"•  {b}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font


def _heading(slide, text, *, color=NAVY):
    _textbox(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.8), text,
             size=30, color=color, bold=True)


# --------------------------------------------------------------- slide types

def add_title_slide(prs, title, subtitle, meta_line):
    """Dark navy title slide. `title` may contain '\\n' for a two-line title."""
    slide = _blank_slide(prs, NAVY)
    _rounded_box(slide, Inches(0.6), Inches(2.55), Inches(0.09), Inches(1.6), AMBER, radius=0.5)
    _textbox(slide, Inches(1.0), Inches(2.3), Inches(11.5), Inches(2.1), title.split("\n"),
             size=44, color=WHITE, bold=True, line_spacing=1.05)
    _textbox(slide, Inches(1.0), Inches(4.5), Inches(11.5), Inches(0.6), subtitle,
             size=18, color=AMBER, bold=True)
    _textbox(slide, Inches(1.0), Inches(6.7), Inches(11.5), Inches(0.5), meta_line,
             size=13, color=GRAY)
    return slide


def add_overview_slide(prs, heading, intro, stages):
    """`stages`: list of (label, description) tuples, rendered as a 3-stage arrow chain."""
    slide = _blank_slide(prs, LIGHT_BLUE)
    _heading(slide, heading)
    _textbox(slide, Inches(0.6), Inches(1.15), Inches(12.1), Inches(1.0), intro,
              size=16, color=GRAY, line_spacing=1.15)

    n = len(stages)
    gap = Inches(0.55)
    box_w = (prs.slide_width - Inches(1.2) - gap * (n - 1)) // n
    box_h = Inches(1.9)
    top = Inches(3.3)
    left = Inches(0.6)
    for i, (label, desc) in enumerate(stages):
        x = left + i * (box_w + gap)
        _rounded_box(slide, x, top, box_w, box_h, DEEP_BLUE)
        _textbox(slide, x + Inches(0.2), top + Inches(0.25), box_w - Inches(0.4), Inches(0.5),
                  label, size=19, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, x + Inches(0.2), top + Inches(0.85), box_w - Inches(0.4), Inches(0.9),
                  desc, size=13, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.1)
        if i < n - 1:
            _right_arrow(slide, x + box_w, top + box_h // 2, gap, AMBER)
    return slide


def add_flow_slide(prs, heading, steps, note=None):
    """`steps`: list of (label, description) tuples, vertical chain with down-arrows."""
    slide = _blank_slide(prs, LIGHT_BLUE)
    _heading(slide, heading)

    n = len(steps)
    box_w = Inches(8.6)
    box_h = Inches(0.85)
    arrow_gap = Inches(0.35)
    left = (prs.slide_width - box_w) // 2
    top = Inches(1.35)
    step_h = box_h + arrow_gap
    for i, (label, desc) in enumerate(steps):
        y = top + i * step_h
        _rounded_box(slide, left, y, box_w, box_h, DEEP_BLUE)
        _textbox(slide, left + Inches(0.3), y + Inches(0.08), box_w - Inches(0.6), Inches(0.35),
                  label, size=15, color=AMBER, bold=True)
        _textbox(slide, left + Inches(0.3), y + Inches(0.42), box_w - Inches(0.6), Inches(0.4),
                  desc, size=12, color=WHITE)
        if i < n - 1:
            _down_arrow(slide, left + box_w // 2, y + box_h + Inches(0.03), Inches(0.3), AMBER)
    if note:
        note_top = top + n * step_h + Inches(0.15)
        _textbox(slide, left, note_top, box_w, Inches(0.6), note,
                  size=12, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    return slide


def add_built_slide(prs, heading, cards):
    """`cards`: list of (subheading, [bullets]) tuples, stacked as cards."""
    slide = _blank_slide(prs, LIGHT_BLUE)
    _heading(slide, heading)

    n = len(cards)
    top = Inches(1.35)
    gap = Inches(0.25)
    card_h = (prs.slide_height - top - Inches(0.4) - gap * (n - 1)) // n
    left = Inches(0.6)
    width = Inches(12.1)
    for i, (subheading, bullets) in enumerate(cards):
        y = top + i * (card_h + gap)
        _rounded_box(slide, left, y, width, card_h, WHITE)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, Inches(0.09), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = AMBER
        bar.line.fill.background()
        bar.shadow.inherit = False
        _textbox(slide, left + Inches(0.35), y + Inches(0.15), width - Inches(0.7), Inches(0.4),
                  subheading, size=16, color=DEEP_BLUE, bold=True)
        tb = slide.shapes.add_textbox(left + Inches(0.35), y + Inches(0.6),
                                       width - Inches(0.7), card_h - Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        _bullets_into(tf, bullets, size=13, color=NAVY, space_after=3)
    return slide


def add_stats_slide(prs, heading, big_number, big_label, col1_heading, col1_bullets,
                     col2_heading, col2_bullets):
    slide = _blank_slide(prs, NAVY)
    _heading(slide, heading, color=WHITE)

    _textbox(slide, Inches(0.6), Inches(1.3), Inches(4.2), Inches(1.6), big_number,
              size=80, color=AMBER, bold=True)
    _textbox(slide, Inches(0.6), Inches(2.85), Inches(4.2), Inches(0.9), big_label,
              size=15, color=WHITE, line_spacing=1.15)

    col_top = Inches(1.35)
    col_w = Inches(3.9)
    col1_left = Inches(5.2)
    col2_left = Inches(9.2)
    for col_left, col_heading, bullets in (
        (col1_left, col1_heading, col1_bullets),
        (col2_left, col2_heading, col2_bullets),
    ):
        _textbox(slide, col_left, col_top, col_w, Inches(0.4), col_heading,
                  size=15, color=AMBER, bold=True)
        tb = slide.shapes.add_textbox(col_left, col_top + Inches(0.5), col_w, Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        _bullets_into(tf, bullets, size=12, color=WHITE, space_after=8)
    return slide


def add_next_steps_slide(prs, heading, items):
    """`items`: list of (title, description) tuples, rendered as a 2-column grid."""
    slide = _blank_slide(prs, LIGHT_BLUE)
    _heading(slide, heading)

    cols = 2
    rows = (len(items) + cols - 1) // cols
    gap = Inches(0.3)
    top = Inches(1.35)
    card_w = (prs.slide_width - Inches(1.2) - gap) // cols
    card_h = (prs.slide_height - top - Inches(0.4) - gap * (rows - 1)) // rows
    left0 = Inches(0.6)
    for i, (title, desc) in enumerate(items):
        r, c = divmod(i, cols)
        x = left0 + c * (card_w + gap)
        y = top + r * (card_h + gap)
        _rounded_box(slide, x, y, card_w, card_h, AMBER_TINT)
        _textbox(slide, x + Inches(0.3), y + Inches(0.2), card_w - Inches(0.6), Inches(0.5),
                  title, size=16, color=DEEP_BLUE, bold=True)
        _textbox(slide, x + Inches(0.3), y + Inches(0.75), card_w - Inches(0.6),
                  card_h - Inches(0.95), desc, size=13, color=NAVY, line_spacing=1.15)
    return slide
